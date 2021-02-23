#%%
from pptx import Presentation

prs = Presentation()
print("Creating new Presentation")
title_slide_layout = prs.slide_layouts[0]

print("Creating slide #1 - with S1-TC name")
slide = prs.slides.add_slide(title_slide_layout)
slide.name = "S1-TC"
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

print("Creating slide #2 - no name")
s2 = prs.slides.add_slide(title_slide_layout)
title = s2.shapes.title
subtitle = s2.placeholders[1]
title.text = "Slide2 Title"
subtitle.text = "No explicit slide Name!"

print("Saving to test.pptx")
prs.save('test.pptx')

#%% Check the ID
newPPT = Presentation("test.pptx")
for slide in newPPT.slides:
    name = slide.name or "#Empty#"
    print(f"Slide name:{name} \tslide_id:{slide.slide_id}")

