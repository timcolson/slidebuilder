#%%
from pptx import Presentation

filename = "source1"
print(f"Read in powerpoint: {filename}.pptx")
ppt = Presentation("../" + filename + ".pptx")

print(ppt)
# %%
for slide in ppt.slides:
    print(f"Slide name:{slide.name}; slide_id:{slide.slide_id}")
# %%

ppt.slides[1].name="TC-NewID-2"

print("Saving file...to " + filename )
ppt.save("./" + filename + "-pyout2.pptx")


# %%
